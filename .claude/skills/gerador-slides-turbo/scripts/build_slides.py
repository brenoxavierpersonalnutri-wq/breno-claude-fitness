#!/usr/bin/env python3
"""
build_slides.py — Pipeline padrão do gerador-slides-turbo.

HTML/CSS (1920x1080) → Chrome headless PNG → PPTX full-bleed (13.333x7.5in).

Por que assim: o desenho nativo do python-pptx (textboxes/shapes) produz slides
visualmente pobres. Renderizar cada slide como HTML dá gradientes, glow radial,
sombras, glassmorphism e tipografia de qualidade — e o PPTX vira só o "envelope"
com os PNGs inseridos full-bleed.

Uso: copie este script, troque os slides de exemplo pela copy real da aula
(recebida do @copywriter-turbo), ajuste WORK/OUT/BRAND e rode.
Cada slide é um bloco HTML dentro de .stage; o CSS compartilhado é o design
system Turbo. Validado em 2026-08-03 (aula Funil 8 do Zero).
"""
import os
import subprocess

from pptx import Presentation
from pptx.util import Inches

CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
WORK = os.path.expanduser("~/tmp/slides-build")   # PNGs/HTMLs intermediários
OUT = os.path.expanduser("~/Desktop/slides.pptx")  # PPTX final
BRAND = "TURBO ACADEMY · NOME DA AULA"             # texto fixo do footer
os.makedirs(WORK, exist_ok=True)

ACCENTS = {"green": "#4ade80", "yellow": "#facc15", "cyan": "#22d3ee",
           "purple": "#a78bfa", "red": "#f87171", "white": "#ffffff"}

# ── Design system Turbo (compartilhado por todos os slides) ──────────────────
CSS = """
*{margin:0;padding:0;box-sizing:border-box}
html,body{width:1920px;height:1080px;overflow:hidden}
body{background:#0a0a0a;color:#fff;position:relative;
  font-family:'Helvetica Neue',-apple-system,'Segoe UI',Arial,sans-serif;
  -webkit-font-smoothing:antialiased}
.glow{position:absolute;width:1400px;height:1400px;border-radius:50%;
  filter:blur(90px);opacity:.13;pointer-events:none}
.glow.tr{top:-700px;right:-500px}.glow.bl{bottom:-800px;left:-600px;opacity:.08}
.stage{position:absolute;inset:0;padding:110px 130px 130px;display:flex;
  flex-direction:column;justify-content:center}
.kicker{font-family:'SF Mono',Menlo,monospace;font-size:22px;letter-spacing:.42em;
  text-transform:uppercase;font-weight:600;margin-bottom:34px}
h1{font-size:92px;line-height:1.04;letter-spacing:-.025em;font-weight:700}
h1 .dim{color:#71717a}
.sub{font-size:32px;color:#a1a1aa;line-height:1.45;margin-top:34px;font-weight:400}
.card{background:rgba(255,255,255,.045);border:1px solid rgba(255,255,255,.09);
  border-radius:26px;padding:52px 58px;position:relative;overflow:hidden;
  box-shadow:0 30px 80px rgba(0,0,0,.5)}
.card::before{content:'';position:absolute;top:0;left:0;right:0;height:5px;
  background:var(--acc)}
.footer{position:absolute;left:130px;right:130px;bottom:56px;display:flex;
  justify-content:space-between;font-family:'SF Mono',Menlo,monospace;
  font-size:19px;color:#52525b;letter-spacing:.24em}
.rule{width:130px;height:6px;background:var(--acc);border-radius:3px;margin-bottom:44px}
ul.b{list-style:none}
ul.b li{font-size:34px;line-height:1.4;color:#d4d4d8;padding-left:52px;
  position:relative;margin-bottom:34px}
ul.b li:last-child{margin-bottom:0}
ul.b li::before{content:'';position:absolute;left:0;top:.52em;width:16px;height:16px;
  border-radius:5px;background:var(--acc)}
.mono{font-family:'SF Mono',Menlo,monospace}
strong{color:#fff;font-weight:700}
.center{align-items:center;text-align:center}
.cols{display:flex;gap:44px}
.cols>.card{flex:1}
"""


def page(body, accent, glow2=None):
    """Envelopa o corpo do slide no HTML completo com glow radial do accent."""
    acc = ACCENTS[accent]
    g2 = f'<div class="glow bl" style="background:{ACCENTS[glow2]}"></div>' if glow2 else ""
    return (f'<!doctype html><html><head><meta charset="utf-8"><style>{CSS}'
            f':root{{--acc:{acc}}}</style></head><body>'
            f'<div class="glow tr" style="background:{acc}"></div>{g2}{body}</body></html>')


S = []  # (body, accent, glow2, com_footer)


def add(body, accent, glow2=None, footer=True):
    S.append((body, accent, glow2, footer))


# ── Slides de exemplo: um por tipo. Substitua pela copy real da aula. ────────

# CAPA (sem footer)
add("""<div class="stage">
<div class="kicker" style="color:#4ade80">KICKER · CONTEXTO DA AULA</div>
<h1>Nome da<br>Aula<span class="dim">.</span></h1>
<div class="sub" style="max-width:1100px">Subtítulo com a promessa da aula,<br>destacando a <strong>palavra-chave</strong>.</div>
<div style="margin-top:70px;font-family:'SF Mono',Menlo,monospace;font-size:24px;color:#4ade80;letter-spacing:.3em">NOME DO EXPERT — MARCA</div>
</div>""", "green", footer=False)

# SECTION DIVIDER (número gigante mono + título)
add("""<div class="stage">
<div class="mono" style="font-size:200px;font-weight:700;color:#f87171;opacity:.92;line-height:1">01</div>
<h1 style="font-size:80px;margin-top:30px">Título da seção</h1>
<div class="sub">Uma linha de contexto sobre o bloco</div>
</div>""", "red")

# CONTEÚDO (card glassmorphism + bullets, máx 4)
add("""<div class="stage">
<div class="kicker" style="color:#4ade80">KICKER DO SLIDE</div>
<h1 style="font-size:68px;margin-bottom:56px">Título do slide:</h1>
<div class="card"><ul class="b">
<li>Primeiro ponto, com <strong>destaque</strong> no que importa</li>
<li>Segundo ponto — frase curta, não parágrafo</li>
<li>Terceiro ponto</li>
</ul></div></div>""", "green")

# COMPARAÇÃO (2 cards, ✗ vermelho vs ✓ verde, glow duplo)
add("""<div class="stage">
<div class="kicker" style="color:#f87171">PERGUNTA QUE POLARIZA</div>
<div class="cols" style="margin-top:26px">
<div class="card" style="--acc:#f87171"><div style="font-size:40px;font-weight:700;color:#f87171;margin-bottom:36px">✗&ensp;O jeito errado</div>
<ul class="b" style="--acc:#f87171"><li style="font-size:31px">sintoma um</li><li style="font-size:31px">sintoma dois</li></ul></div>
<div class="card"><div style="font-size:40px;font-weight:700;color:#4ade80;margin-bottom:36px">✓&ensp;O jeito certo</div>
<ul class="b"><li style="font-size:31px">contraste um</li><li style="font-size:31px">contraste dois</li></ul></div>
</div></div>""", "green", glow2="red")

# QUOTE / DESTAQUE (aspas Georgia gigantes)
add("""<div class="stage center">
<div style="font-family:Georgia,serif;font-size:190px;color:#4ade80;line-height:.4;margin-bottom:48px">“</div>
<h1 style="font-size:70px;max-width:1560px;font-weight:600">A frase de impacto da aula,<br>com o <span style="color:#4ade80">trecho-chave</span> no accent.</h1>
</div>""", "green")

# NÚMERO / DADO (mono gigante)
add("""<div class="stage center">
<div class="mono" style="font-size:290px;font-weight:700;color:#4ade80;letter-spacing:-.03em;line-height:1">R$ 0</div>
<div class="sub" style="font-size:38px;color:#fff;max-width:1400px;margin-top:40px">o label que explica<br>o que esse número significa</div>
</div>""", "green")

# ENCERRAMENTO (sem footer)
add("""<div class="stage center">
<h1 style="font-size:96px;letter-spacing:.02em">MARCA<span style="color:#4ade80">.</span></h1>
<div class="sub" style="font-size:30px">Frase de fechamento / transição pro Q&A.</div>
</div>""", "green", footer=False)


# ── Render: HTML → PNG (Chrome headless) → PPTX full-bleed ───────────────────
total = len(S)
pngs = []
for i, (body, accent, glow2, com_footer) in enumerate(S, 1):
    if com_footer:
        body += (f'<div class="footer"><span>{BRAND}</span>'
                 f'<span>{i:02d} / {total}</span></div>')
    hp = os.path.join(WORK, f"s{i:02d}.html")
    pp = os.path.join(WORK, f"s{i:02d}.png")
    with open(hp, "w") as f:
        f.write(page(body, accent, glow2))
    subprocess.run([CHROME, "--headless", "--disable-gpu", "--hide-scrollbars",
                    "--force-device-scale-factor=1", f"--screenshot={pp}",
                    "--window-size=1920,1080", f"file://{hp}"],
                   capture_output=True, timeout=60)
    if not os.path.exists(pp):
        raise SystemExit(f"screenshot falhou no slide {i}")
    pngs.append(pp)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
blank = prs.slide_layouts[6]
for pp in pngs:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(pp, 0, 0, width=Inches(13.333), height=Inches(7.5))
prs.save(OUT)
print(f"OK: {len(pngs)} slides → {OUT}")
